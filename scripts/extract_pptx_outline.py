"""
Extraction hors-ligne du plan du PPTX "Pharmacométrie Pratique.pptx".

Usage (dans ce repo) :
    python scripts/extract_pptx_outline.py Pharmacométrie\ Pratique.pptx src/content/pptx_outline.json

Prérequis :
    pip install python-pptx

Ce script n'est pas appelé par le site en production ; il sert uniquement
à générer un plan JSON exploitable pour enrichir le contenu.
"""

import json
import sys
from pptx import Presentation


def slide_dict(slide, idx):
    title = slide.shapes.title.text if slide.shapes.title else f"Slide {idx+1}"
    bullets = []
    notes = ''
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text.strip()
        if text and text != title:
            bullets.append(text)
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()
    return {"index": idx + 1, "title": title, "bullets": bullets, "notes": notes}


def main():
    if len(sys.argv) < 3:
        print("Usage: extract_pptx_outline.py input.pptx output.json")
        sys.exit(1)
    input_path, output_path = sys.argv[1], sys.argv[2]
    prs = Presentation(input_path)
    slides = [slide_dict(s, i) for i, s in enumerate(prs.slides)]
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump({"slides": slides}, f, ensure_ascii=False, indent=2)
    print(f"Wrote {len(slides)} slides to {output_path}")


if __name__ == "__main__":
    main()
